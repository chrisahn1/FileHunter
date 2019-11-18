import xlsxwriter
from xlsxwriter.utility import xl_rowcol_to_cell
import glob
import xlrd
import tika
import time
import string
import os
import re
import sys
import glob
import PyPDF2 # pip install PyPDF2
import docx # pip install python-docx
import PyQt5 # pip install pyqt5
from os import walk
import mmap
from tika import parser
from threading import  Thread, Lock
from datetime import datetime
import textract
import pickle
import multiprocessing as mp
from math import cos
from pptx import Presentation
mutex = Lock() # this the lock to prevent race condition
drivesList = [] # list of Drives available on the machine.
drivesFolders = [] # list of all root folders in the computers
processList = [] # the list that contains all the threads
CandidateList = [] # the lsit of all files in the machine that ends with the specified extension.
domainSet = [] # Refined version of CandidateList but with all the constraints applied on the CandidateList.
ConstraintsList = [] # a list of all no folders that we are outside of our interest. No search will be conducted on those folders or locations.
extension_List = [] # a list of desired extensions. The extensions specified here will be applied to obtain the CandidateList.
MatchList = [] #the list of files where a match to the keyword is found. This is the lsit that will be returned to the GUI

###########################################
# extension_List = [".txt", ".doc", ".docx", ".wpd", ".wp", ".wp7", ".dto", ".dotx"
#                   ".ppt", ".pptx" ,
#                   ".xls", ".csv",
#                     ".pdf"]

def getExtensionList():
    extension_List = [ ".ppt", ".pptx", ".txt", ".doc", ".docx", "xls", "xlsx"]
    return extension_List

def getConstraintsList ():
    ConstraintsList = ['C:/FILESDB','C:/Windows','C:/Program Files','C:/Program Files (x86)','C:/pagefile.sys','C:/swapfile.sys','C:/Lib','C:/libs','C:/$Recycle.Bin', 'Local']
    return ConstraintsList

def getPathRestrictions():
    pathRestrictionsList = ['lib', 'sys', 'Sys', 'SYS', 'Lib', 'AppData', 'embedded', '~$']
    return pathRestrictionsList

t1 = datetime.now()
os.chdir('/')
def get_drives():
	response = os.popen("wmic logicaldisk get caption")
	list1 = []
	total_file = []

	for line in response.readlines():
		line = line.strip("\n")
		line = line.strip("\r")
		line = line.strip(" ")
		if (line == "Caption" or line == ""):
			continue
		list1.append(line)
	return list1

def getDrivesFolders(drivesList):
    drivesFolders = []
    for drive in drivesList:
        try:
            driveContent = os.listdir(drive)
            for item in driveContent:
                pathedItem = drive +'/'+item
                drivesFolders.append(pathedItem)
        except:
            continue
    return drivesFolders
        #I used try / except because sometimes reading everydrive is not possible. for example, DVD drive can only be read if there is a DVD inside in the drive.
        #Therefore, we are ignoring the drives that has no data.

def removeSysFolders(driversFolder):
    refinedList=[]
    sysFoldersList = getConstraintsList ()
    for item in driversFolder:
        try:
            if  item in sysFoldersList:
                continue
            else:
                refinedList.append(item)
        except:
                continue
    return refinedList

def spider(item):
    global CandidateList
    constraintList = getConstraintsList()
    foundConstraints = False
    extension_List = getExtensionList()
    for extension in extension_List:
        if item.lower().endswith(extension):
            CandidateList.append(item)
            return
    for (dirname,dirs,files) in os.walk(item):
                for filename in files:
                    for constraint in constraintList:
                        if (filename.find(constraint) > -1):
                            foundConstraints = True
                    if (~foundConstraints):
                        for extension in extension_List:
                            if dirname.lower().endswith(extension):
                                CandidateList.append(dirname)
                            elif filename.lower().endswith(extension):
                                cwd = os.getcwd()
                                fullPath  = os.path.join(dirname,filename)
                                CandidateList.append(fullPath)
							#print(fullPath)
							#if os.path.isfile(fullPath):
								#print(fullPath)
								#result.append(fullPath)
########

#spider(folderPath)
def refineList():
    global CandidateList
    newList = []
    pathResList = getPathRestrictions()
    #mutex.acquire()
    #try:
    for file in CandidateList:
            found = False
            for constrain in getPathRestrictions():
                if(file.find(constrain) > -1):
                #if re.search(file.lower(), constrain.lower()):
                    #try:
                        found = True
                    #except:
                        #continue
            if (found == False):
                newList.append(file)
    #finally:
        #mutex.release()
    return newList

def creat():
    drivesList = get_drives()
    drivesFolders = getDrivesFolders(drivesList)
    refinedList = removeSysFolders((drivesFolders))
    print(refinedList)
    processList = []
    mutex.acquire()
    try:
        for item in refinedList:
            process1 = Thread(target = spider, args =(item,))
            process1.start()
            processList.append(process1)
        for t in processList:
            t.join()
        print(processList)
    finally:
        mutex.release()
    domainSet1 = refineList()
    makePickleFile(domainSet1)


def searchInPDF(filename, key):
    occurrences = 0
    pdfFileObj = open(filename,'rb')
    pdfReader = PyPDF2.PdfFileReader(pdfFileObj)
    num_pages = pdfReader.numPages
    count = 0
    text = ""
    while count < num_pages:
        pageObj = pdfReader.getPage(count)
        count +=1
        text += pageObj.extractText()
    if text != "":
       text = text
    else:
       text = textract.process(filename, method='tesseract', language='eng')
    tokens = word_tokenize(text)
    punctuation = ['(',')',';',':','[',']',',']
    stop_words = stopwords.words('english')
    keywords = [word for word in tokens if not word in stop_words and  not word in punctuation]
    for k in keywords:
        if key == k: occurrences+=1
    return occurrences


def inList(list, item):
    for e in list:
        if e == item:
            return True
    return False

def binaryFile(file, keyword):
    with open(file)as f:
        s = mmap.mmap(f.fileno(), 0, access=mmap.ACCESS_READ)
        if s.find(keyword) != -1 :
            return True
    return False

def searchBinaryFile(file, keyword):
    s= inFile(file, 'rb').read()
    return (s.find(keyword) >-1)

def searchDomain( keyword):
    domainSet = getPickleFile()
    MatchList = []
    for file in domainSet:
        if (file_size(file) < float(0x19000)):  # maximum file size to read is 100 kilo bytes.
            resSearch = ""
            try:
                if(searchFiles(file, keyword)):
                    MatchList.append(file)
            except:
                #print('Failed to read this file ' + file)
                continue
    if (len(MatchList) == 0):
        MatchList.append("SORRY...NO RESULTS ARE FOUND")
    return MatchList

def searchTxtFile(file, keyword):
    found = False
    if (file.lower().endswith('txt')):
        # f = open(file, 'r')
        with open(file) as fp:
            for cnt, line in enumerate(fp):
                if (line.find(keyword) > -1):
                    #print('FOUUUUUUUUUUUUUUUUUND')
                    # print("line{}: {}".format(cnt, line))
                    found = True
    return found

def searchFiles(file, keyword):
    found = False
    formatFile = file.replace('\\','\\\\')
    #print(formatFile)
    if (file.lower().endswith('.ppt') or file.lower().endswith('.pptx')):
        found = readPPTFiles(file, keyword)
        return found
    elif (file.lower().endswith('.doc') or file.lower().endswith('.docx') ):
        found = readDocx(formatFile, keyword)
        return found
    elif  (file.lower().endswith('.txt')):
        found = searchTxtFile(formatFile, keyword)
        return found
    elif (file.lower().endswith('.xlsx') or file.lower().endswith('.xls')):
        found = readEXCELFiles(formatFile, keyword)
        return found
    else:
        return found




def readDocx(file, keyword):
    doc = docx.Document(file)
    num = len( doc.paragraphs)
    text= []
    for i in range(0, num):
        text.append(doc.paragraphs[i].text)
    #print(text)
    for word in text:
        if word.find(keyword) != -1:
            return True
    return False

def readPDF(file):
    doc = textract.process(file)
    print(doc)

def readPPTFiles(file, keyword):
        #print('Hello from the readPPTFiles')
        formatFile = file.replace('\\', '\\\\')
        prs = Presentation(formatFile)
        #print(formatFile)
        #print("----------------------")
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    #print(shape.text)
                    if(shape.text.find(keyword)> -1):
                        print('----**//**//**/ FOUND')
                        return True
        return False

def makePickleFile(searchDomain):
    filename = r'C:\Users\Public\DataDB'
    outFile = open(filename, "wb")
    pickle.dump(searchDomain, outFile)
    outFile.close()


def getPickleFile():
    filename = r'C:\Users\Public\DataDB'
    inFile = open(filename,'rb')
    domainSet = pickle.load(inFile)
    #print(searchDomain)
    inFile.close()
    return domainSet
#

def convert_bytes(num):
    numFloat = float(num)
    """
    this function will convert bytes to MB.... GB... etc
    """
    # for x in ['bytes', 'KB', 'MB', 'GB', 'TB']:
    #     if num < 1024.0:
            # return "%3.1f %s" % (num, x)
    return numFloat
        # num /= 1024.0


def file_size(file_path):
    """
    this function will return the file size
    """
    if os.path.isfile(file_path):
        file_info = os.stat(file_path)
        return convert_bytes(file_info.st_size)


# makePickleFile(domainSet)
# for i in CandidateList:
#     print(i)
# print(len(CandidateList))

# time.sleep(10)

#domainSet = getPickleFile()
# # print(len(searchDomain))

# this file will be talking two parameters
# 1. a file path for a excel file
# 2. the key word that will be searched in that file
# this class will return only true or false


#test purpose
#path = (r"C:\Users\17148\PycharmProjects\exceltester\testers\yeah.xls")
#term = "abc"

def readEXCELFiles(file, keyword):
    # it can be poosbile that the passing path is not availabel for this function
    # if that happens
    # please try to convert the path using
    # path = (r"the_path_of_the_file")
    # and then call this function
    # everything should be fine
    if (file_size(file) < float(0xD0900)): #maximum file size to read is five mega bytes.

        for sh in xlrd.open_workbook(file).sheets():

                for row in range(sh.nrows):
                    for col in range(sh.ncols):
                        myCell = sh.cell(row, col)
                        #print(myCell)
                        if myCell.value == keyword:
                            #print(myCell)
                    # once the term is found
                    # immediatly return true
                    # and stop the loop
                    # for better efficiency
                            return True
                            break


    return False


#domainSet = getPickleFile()
# # print(len(searchDomain))
# creat()
#keyword="The security of a system, application, or protocol is always relative"

# domainSet = refineList()
# for i in domainSet:
#     print(i, file_size(i))
# makePickleFile(domainSet)
# domainset = getPickleFile()
# print(len(domainSet))
# #
# print("-----------########----------------")
# keyword = "Ahmed"
# domSet = getPickleFile()
# MatchList = searchDomain( keyword)
# for i in domSet:
#     print(i)
#

t2 = datetime.now()
totalTime = t2-t1
print(totalTime)
